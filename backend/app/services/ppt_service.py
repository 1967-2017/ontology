from __future__ import annotations

from io import BytesIO
from pathlib import Path
import re
from typing import Any

from fastapi import HTTPException
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from sqlalchemy.orm import Session

from app.config import get_settings
from app.db.models import PresentationModel, PresentationStatus
from app.services.document_service import DocumentService
from app.services.knowledge_service import KnowledgeService
from app.services.ocr_service import OCRService
from app.services.ppt_image_service import PPTImageService
from app.services.ppt_planner_service import PPTPlannerService


class PPTService:
    TARGET_SLIDE_SIZE = (1600, 900)
    TARGET_ASPECT_RATIO = 16 / 9
    ASPECT_RATIO_TOLERANCE = 0.01

    def __init__(
        self,
        knowledge_service: KnowledgeService,
        document_service: DocumentService | None = None,
        ocr_service: OCRService | None = None,
        planner_service: PPTPlannerService | None = None,
        image_service: PPTImageService | None = None,
    ) -> None:
        settings = get_settings()
        self.knowledge_service = knowledge_service
        self.document_service = document_service or DocumentService()
        self.ocr_service = ocr_service or OCRService()
        self.planner_service = planner_service or PPTPlannerService()
        self.image_service = image_service or PPTImageService()
        self.base_dir = Path(settings.data_dir)
        self.export_dir = self.base_dir / settings.export_dir_name
        self.export_dir.mkdir(parents=True, exist_ok=True)

    def generate(
        self,
        db: Session,
        topic: str | None,
        slide_count: int,
        document_ids: list[int] | None,
        use_knowledge_base: bool,
    ) -> dict[str, Any]:
        resolved_topic = (topic or "").strip()
        if not resolved_topic and not use_knowledge_base and not document_ids:
            raise HTTPException(status_code=400, detail="生成 PPTX 汇报需要主题、已上传文档或知识库资料")

        document_labels = self._resolve_document_labels(db, document_ids or [])
        prompt = resolved_topic or ("、".join(document_labels[:2]) if document_labels else "上传资料汇报")
        context_text = self._build_context_text(db, prompt, document_ids, use_knowledge_base)
        if not context_text.strip():
            raise HTTPException(status_code=400, detail="未能提取到可用于生成 PPTX 汇报的文本内容")

        outline_data = self.planner_service.generate_outline(prompt, slide_count, context_text)
        slides = self._normalize_outline_slides(outline_data.get("slides", []), prompt)
        if not slides:
            raise HTTPException(status_code=500, detail="大纲生成成功但未返回任何页面")

        presentation = PresentationModel(
            title=str(slides[0].get("title") or prompt),
            topic=prompt,
            status=PresentationStatus.pending,
            slide_count=slide_count,
            source_document_ids=document_ids,
            outline=slides,
        )
        db.add(presentation)
        db.commit()
        db.refresh(presentation)

        slide_images = self._render_slide_images(slides, presentation.id)
        output_path = self.export_dir / f"presentation-{presentation.id}.pptx"
        self._render_pptx(output_path, slide_images)

        presentation.file_path = str(output_path)
        presentation.status = PresentationStatus.completed
        db.commit()
        if document_ids:
            self._consume_source_documents(db, document_ids)
        db.refresh(presentation)
        return self.serialize_presentation(presentation)

    def get_presentation(self, db: Session, presentation_id: int) -> dict[str, Any]:
        presentation = db.get(PresentationModel, presentation_id)
        if presentation is None:
            raise HTTPException(status_code=404, detail="PPT 任务不存在")
        return self.serialize_presentation(presentation)

    def serialize_presentation(self, presentation: PresentationModel) -> dict[str, Any]:
        return {
            "presentation_id": presentation.id,
            "title": presentation.title,
            "topic": presentation.topic,
            "status": presentation.status.value,
            "slide_count": presentation.slide_count,
            "download_url": f"/api/backend/ppt/{presentation.id}/download" if presentation.file_path else None,
            "outline": presentation.outline,
            "source_document_ids": presentation.source_document_ids,
        }

    def _normalize_outline_slides(self, slides: list[dict[str, Any]], prompt: str) -> list[dict[str, Any]]:
        if not slides:
            return []
        normalized = [dict(slide) for slide in slides]
        normalized[0]["title"] = str(normalized[0].get("title") or prompt)
        normalized[-1]["title"] = "结论与行动建议"
        if not normalized[-1].get("content_summary"):
            normalized[-1]["content_summary"] = "总结核心判断并给出后续行动建议"
        if not normalized[-1].get("visual_prompt"):
            normalized[-1]["visual_prompt"] = "A clean closing presentation slide for conclusions and action recommendations."
        for index, slide in enumerate(normalized):
            slide["index"] = int(slide.get("index", index))
        return normalized

    def _build_context_text(
        self,
        db: Session,
        prompt: str,
        document_ids: list[int] | None,
        use_knowledge_base: bool,
    ) -> str:
        if use_knowledge_base:
            answer = self.knowledge_service.answer(db, prompt, limit=8)
            return self._summarize_context([citation["snippet"] for citation in answer["citations"]], prompt)
        if document_ids:
            snippets = self._collect_uploaded_document_material(db, document_ids, 18)
            return self._summarize_context(snippets, prompt)
        return prompt

    def _render_slide_images(self, slides: list[dict[str, Any]], presentation_id: int) -> list[Path]:
        rendered_paths: list[Path] = []
        previous_prompt: str | None = None
        for slide in slides:
            visual_prompt = str(slide.get("visual_prompt") or slide.get("content_summary") or slide.get("title") or "")
            image_bytes = self.image_service.generate_slide_image(visual_prompt, previous_prompt)
            self._validate_slide_image(image_bytes)
            image_path = self.export_dir / f"presentation-{presentation_id}-slide-{slide['index']}.png"
            image_path.write_bytes(image_bytes)
            rendered_paths.append(image_path)
            previous_prompt = visual_prompt
        return rendered_paths

    def _validate_slide_image(self, image_bytes: bytes) -> None:
        with Image.open(BytesIO(image_bytes)) as source_image:
            width, height = source_image.size
        if width <= 0 or height <= 0:
            raise HTTPException(status_code=500, detail="图片模型返回了无效尺寸的页面图片")
        ratio = width / height
        if abs(ratio - self.TARGET_ASPECT_RATIO) > self.ASPECT_RATIO_TOLERANCE:
            raise HTTPException(
                status_code=500,
                detail=(
                    f"图片模型返回的页面图片比例不是 16:9：实际为 {width}x{height}。"
                    "当前已禁用服务端裁切和补边，请让上游直接输出 16:9 图片。"
                ),
            )

    def _render_pptx(self, output_path: Path, slide_images: list[Path]) -> None:
        if not slide_images:
            raise HTTPException(status_code=500, detail="没有可用于输出 PPTX 的页面图片")
        presentation = Presentation()
        presentation.slide_width = Inches(13.333333)
        presentation.slide_height = Inches(7.5)
        blank_layout = presentation.slide_layouts[6]
        for slide_path in slide_images:
            slide = presentation.slides.add_slide(blank_layout)
            slide.shapes.add_picture(str(slide_path), 0, 0, width=presentation.slide_width, height=presentation.slide_height)
        presentation.save(output_path)

    def _resolve_document_labels(self, db: Session, document_ids: list[int]) -> list[str]:
        labels: list[str] = []
        for document_id in document_ids:
            document = self.document_service.get_document_model(db, document_id)
            labels.append(document.filename)
        return labels

    def _consume_source_documents(self, db: Session, document_ids: list[int]) -> None:
        for document_id in document_ids:
            document = self.document_service.get_document_model(db, document_id)
            if document.workspace_active:
                self.document_service.mark_document_consumed(db, document)

    def _collect_uploaded_document_material(self, db: Session, document_ids: list[int], limit: int) -> list[str]:
        snippets: list[str] = []
        for document_id in document_ids:
            document = self.document_service.get_document_model(db, document_id)
            ocr_result = self.document_service.get_ocr_result_model(db, document_id)
            if ocr_result is not None and ocr_result.full_text.strip():
                snippets.extend(self._to_snippets(ocr_result.full_text, limit - len(snippets)))
            else:
                extracted = self.ocr_service.extract(document.storage_path, document.file_type)
                snippets.extend(self._page_text_to_snippets(extracted.get("pages", []), limit - len(snippets)))
            if len(snippets) >= limit:
                break
        return snippets[:limit]

    def _page_text_to_snippets(self, pages: list[dict[str, Any]], remaining: int) -> list[str]:
        collected: list[str] = []
        for page in pages:
            text = str(page.get("text") or "").strip()
            if not text:
                continue
            collected.extend(self._to_snippets(text, remaining - len(collected)))
            if len(collected) >= remaining:
                break
        return collected[:remaining]

    def _to_snippets(self, text: str, limit: int) -> list[str]:
        if limit <= 0:
            return []
        parts = [part.strip() for part in re.split(r"\n{2,}|(?<=[。！？])", text.replace("\r", "\n")) if part.strip()]
        if not parts:
            parts = [text.strip()]
        return parts[:limit]

    def _summarize_context(self, snippets: list[str], prompt: str) -> str:
        cleaned = [snippet.strip() for snippet in snippets if snippet and snippet.strip()]
        if not cleaned:
            return ""
        summary_lines = [f"主题：{prompt}", "资料摘要："]
        for index, snippet in enumerate(cleaned[:10], start=1):
            summary_lines.append(f"{index}. {snippet}")
        return "\n".join(summary_lines)
