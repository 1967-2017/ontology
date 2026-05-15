from pathlib import Path

from app.api import documents as documents_api
from app.api import knowledge as knowledge_api
from app.api import ppt as ppt_api
from app.db.models import DocumentOCRStatus
from pptx import Presentation


def configure_document_environment(tmp_path: Path):
    document_dir = tmp_path / "documents"
    knowledge_dir = tmp_path / "knowledge"
    export_dir = tmp_path / "exports"
    document_dir.mkdir(parents=True, exist_ok=True)
    knowledge_dir.mkdir(parents=True, exist_ok=True)
    export_dir.mkdir(parents=True, exist_ok=True)

    documents_api.document_service.storage_dir = document_dir
    documents_api.knowledge_service.knowledge_dir = knowledge_dir

    knowledge_api.knowledge_service.knowledge_dir = knowledge_dir

    ppt_api.ppt_service.knowledge_service.knowledge_dir = knowledge_dir

    ppt_api.ppt_service.export_dir = export_dir


def test_document_ocr_rag_and_ppt_flow(client, tmp_path, monkeypatch):
    configure_document_environment(tmp_path)

    upload_response = client.post(
        "/documents/upload",
        files={"file": ("scan.pdf", b"%PDF-1.4 mock", "application/pdf")},
    )
    assert upload_response.status_code == 200
    document_id = upload_response.json()["data"]["document_id"]

    monkeypatch.setattr(
        documents_api.ocr_service,
        "extract",
        lambda file_path, file_type: {
            "full_text": "AI Agent 项目规划\n需要先做 OCR 再做知识库。",
            "pages": [
                {
                    "page_number": 1,
                    "text": "AI Agent 项目规划\n需要先做 OCR 再做知识库。",
                    "blocks": [{"text": "AI Agent 项目规划"}],
                    "source": "ocr",
                }
            ],
            "blocks": [{"text": "AI Agent 项目规划"}],
        },
    )

    ocr_response = client.post(f"/documents/{document_id}/ocr")
    assert ocr_response.status_code == 200
    assert "AI Agent" in ocr_response.json()["data"]["full_text"]

    index_response = client.post(f"/documents/{document_id}/index")
    assert index_response.status_code == 200
    assert index_response.json()["data"]["indexed_chunks"] >= 1

    answer_response = client.post("/knowledge/answer", json={"question": "AI Agent"})
    assert answer_response.status_code == 200
    answer_payload = answer_response.json()["data"]
    assert answer_payload["citations"][0]["filename"] == "scan.pdf"
    assert "AI Agent" in answer_payload["answer"]

    monkeypatch.setattr(
        ppt_api.ppt_service.ocr_service,
        "extract",
        lambda file_path, file_type: {
            "full_text": "AI Agent 项目规划\n需要先做 OCR 再做知识库。",
            "pages": [
                {
                    "page_number": 1,
                    "text": "AI Agent 项目规划\n需要先做 OCR 再做知识库。",
                    "blocks": [],
                    "source": "ocr",
                }
            ],
            "blocks": [],
        },
    )
    monkeypatch.setattr(
        ppt_api.ppt_service.planner_service,
        "generate_outline",
        lambda topic, slide_count, context_text: {
            "global_style": "clean tech",
            "slides": [
                {
                    "index": index,
                    "title": f"Slide {index + 1}",
                    "content_summary": f"Summary {index + 1}",
                    "visual_prompt": f"Visual {index + 1}",
                }
                for index in range(slide_count)
            ],
        },
    )
    monkeypatch.setattr(
        ppt_api.ppt_service.image_service,
        "generate_slide_image",
        lambda visual_prompt, reference_style_prompt=None: _build_test_png_bytes(),
    )
    ppt_response = client.post(
        "/ppt/generate",
        json={"topic": "AI Agent", "slide_count": 4, "use_knowledge_base": True},
    )
    assert ppt_response.status_code == 200
    ppt_payload = ppt_response.json()["data"]
    assert ppt_payload["slide_count"] == 4
    assert ppt_payload["download_url"].endswith("/download")

    download_response = client.get(ppt_payload["download_url"].replace("/api/backend", ""))
    assert download_response.status_code == 200
    assert (
        download_response.headers["content-type"]
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert download_response.content[:2] == b"PK"

    delete_response = client.delete(f"/documents/{document_id}")
    assert delete_response.status_code == 200

    answer_after_delete = client.post("/knowledge/answer", json={"question": "AI Agent"})
    assert answer_after_delete.status_code == 200
    assert answer_after_delete.json()["data"]["citations"] == []
    document_list_after_delete = client.get("/documents")
    assert document_list_after_delete.status_code == 200
    assert document_list_after_delete.json()["data"]["items"] == []


def test_document_ocr_sanitizes_invalid_surrogates_before_persisting(client, tmp_path, monkeypatch):
    configure_document_environment(tmp_path)

    upload_response = client.post(
        "/documents/upload",
        files={"file": ("scan.png", b"fake-image", "image/png")},
    )
    document_id = upload_response.json()["data"]["document_id"]

    monkeypatch.setattr(
        documents_api.ocr_service,
        "extract",
        lambda file_path, file_type: {
            "full_text": "alpha\ud835omega",
            "pages": [{"page_number": 1, "text": "beta\ud835gamma", "blocks": [{"text": "delta\ud835theta"}]}],
            "blocks": [{"text": "delta\ud835theta"}],
        },
    )

    ocr_response = client.post(f"/documents/{document_id}/ocr")
    assert ocr_response.status_code == 200
    payload = ocr_response.json()["data"]
    assert payload["full_text"] == "alphaomega"
    assert payload["pages"][0]["text"] == "betagamma"
    assert payload["blocks"][0]["text"] == "deltatheta"


def test_get_document_ocr_result_returns_processing_state_without_rerun(client, db_session, tmp_path, monkeypatch):
    configure_document_environment(tmp_path)

    upload_response = client.post(
        "/documents/upload",
        files={"file": ("scan.png", b"fake-image", "image/png")},
    )
    document_id = upload_response.json()["data"]["document_id"]

    document = documents_api.document_service.get_document_model(db_session, document_id)
    document.ocr_status = DocumentOCRStatus.processing
    db_session.commit()

    monkeypatch.setattr(
        documents_api.ocr_service,
        "extract",
        lambda file_path, file_type: (_ for _ in ()).throw(AssertionError("OCR should not rerun while processing")),
    )

    get_response = client.get(f"/documents/{document_id}/ocr")
    assert get_response.status_code == 200
    assert get_response.json()["data"]["status"] == "processing"

    post_response = client.post(f"/documents/{document_id}/ocr")
    assert post_response.status_code == 200
    payload = post_response.json()["data"]
    assert payload["status"] == "processing"
    assert "处理中" in payload["message"]


def test_run_document_ocr_returns_failed_state_without_rerun(client, db_session, tmp_path, monkeypatch):
    configure_document_environment(tmp_path)

    upload_response = client.post(
        "/documents/upload",
        files={"file": ("scan.png", b"fake-image", "image/png")},
    )
    document_id = upload_response.json()["data"]["document_id"]

    document = documents_api.document_service.get_document_model(db_session, document_id)
    document.ocr_status = DocumentOCRStatus.failed
    db_session.commit()

    monkeypatch.setattr(
        documents_api.ocr_service,
        "extract",
        lambda file_path, file_type: (_ for _ in ()).throw(AssertionError("Failed OCR should not auto-rerun")),
    )

    response = client.post(f"/documents/{document_id}/ocr")
    assert response.status_code == 200
    payload = response.json()["data"]
    assert payload["status"] == "failed"
    assert "失败" in payload["message"]


def test_generate_ppt_from_uploaded_document_without_prior_ocr_or_kb(client, tmp_path, monkeypatch):
    configure_document_environment(tmp_path)

    upload_response = client.post(
        "/documents/upload",
        files={"file": ("brief.pdf", b"%PDF-1.4 mock", "application/pdf")},
    )
    assert upload_response.status_code == 200
    document_id = upload_response.json()["data"]["document_id"]

    monkeypatch.setattr(
        ppt_api.ppt_service.ocr_service,
        "extract",
        lambda file_path, file_type: {
            "full_text": "项目背景\n实施路径\n资源安排\n风险控制",
            "pages": [
                {"page_number": 1, "text": "项目背景\n实施路径", "blocks": [], "source": "ocr"},
                {"page_number": 2, "text": "资源安排\n风险控制", "blocks": [], "source": "ocr"},
            ],
            "blocks": [],
        },
    )
    monkeypatch.setattr(
        ppt_api.ppt_service.planner_service,
        "generate_outline",
        lambda topic, slide_count, context_text: {
            "global_style": "clean tech",
            "slides": [
                {
                    "index": 0,
                    "title": topic,
                    "content_summary": "封面页",
                    "visual_prompt": "封面大图",
                },
                {
                    "index": 1,
                    "title": "背景与问题",
                    "content_summary": "围绕项目背景与问题展开",
                    "visual_prompt": "背景与问题的科技风可视化",
                },
                {
                    "index": 2,
                    "title": "方法与路径",
                    "content_summary": "说明实施路径与方法",
                    "visual_prompt": "方法流程图形化展示",
                },
                {
                    "index": 3,
                    "title": "资源安排",
                    "content_summary": "总结资源安排",
                    "visual_prompt": "资源配置与卡片式布局",
                },
                {
                    "index": 4,
                    "title": "风险控制",
                    "content_summary": "说明风险控制重点",
                    "visual_prompt": "风险与治理的简洁展示",
                },
                {
                    "index": 5,
                    "title": "结论与建议",
                    "content_summary": "收束为行动建议",
                    "visual_prompt": "结论收束页",
                },
            ],
        },
    )
    monkeypatch.setattr(
        ppt_api.ppt_service.image_service,
        "generate_slide_image",
        lambda visual_prompt, reference_style_prompt=None: _build_test_png_bytes(),
    )

    response = client.post(
        "/ppt/generate",
        json={"document_ids": [document_id], "slide_count": 6},
    )
    assert response.status_code == 200
    payload = response.json()["data"]
    assert payload["slide_count"] == 6
    assert payload["topic"] == "brief.pdf"
    assert payload["download_url"].endswith("/download")
    outline = payload["outline"]
    assert outline[0]["title"] == "brief.pdf"
    assert outline[1]["title"] == "背景与问题"
    assert outline[-1]["title"] == "结论与行动建议"
    assert outline[2]["content_summary"] == "说明实施路径与方法"


def test_generate_ppt_from_uploaded_document_fails_when_no_text_can_be_extracted(client, tmp_path, monkeypatch):
    configure_document_environment(tmp_path)

    upload_response = client.post(
        "/documents/upload",
        files={"file": ("empty.png", b"fake-image", "image/png")},
    )
    document_id = upload_response.json()["data"]["document_id"]

    monkeypatch.setattr(
        ppt_api.ppt_service.ocr_service,
        "extract",
        lambda file_path, file_type: {"full_text": "", "pages": [], "blocks": []},
    )

    response = client.post(
        "/ppt/generate",
        json={"document_ids": [document_id], "slide_count": 6},
    )
    assert response.status_code == 400
    assert "未能提取到可用于生成 PPTX 汇报的文本内容" in response.json()["error"]["message"]


def test_ppt_outline_changes_with_slide_count_and_document_structure(client, tmp_path, monkeypatch):
    configure_document_environment(tmp_path)

    upload_response = client.post(
        "/documents/upload",
        files={"file": ("strategy.pdf", b"%PDF-1.4 mock", "application/pdf")},
    )
    document_id = upload_response.json()["data"]["document_id"]

    monkeypatch.setattr(
        ppt_api.ppt_service.ocr_service,
        "extract",
        lambda file_path, file_type: {
            "full_text": (
                "背景：团队协作复杂度上升。\n\n"
                "定义：Vibe Coding 是一种以自然语言驱动的软件协作方式。\n\n"
                "方法：通过提示词、快速试错和人机协同完成开发。\n\n"
                "优势：提升原型效率，缩短验证周期。\n\n"
                "风险：代码质量和边界控制需要额外治理。\n\n"
                "建议：先从低风险场景试点，再沉淀规范。"
            ),
            "pages": [
                {
                    "page_number": 1,
                    "text": (
                        "背景：团队协作复杂度上升。\n\n"
                        "定义：Vibe Coding 是一种以自然语言驱动的软件协作方式。\n\n"
                        "方法：通过提示词、快速试错和人机协同完成开发。\n\n"
                        "优势：提升原型效率，缩短验证周期。\n\n"
                        "风险：代码质量和边界控制需要额外治理。\n\n"
                        "建议：先从低风险场景试点，再沉淀规范。"
                    ),
                    "blocks": [],
                    "source": "ocr",
                }
            ],
            "blocks": [],
        },
    )
    monkeypatch.setattr(
        ppt_api.ppt_service.planner_service,
        "generate_outline",
        lambda topic, slide_count, context_text: {
            "global_style": "clean tech",
            "slides": [
                {"index": 0, "title": topic, "content_summary": "封面", "visual_prompt": "封面"},
                {"index": 1, "title": "背景", "content_summary": "背景", "visual_prompt": "背景"},
                {"index": 2, "title": "概念", "content_summary": "概念", "visual_prompt": "概念"},
                {"index": 3, "title": "方法", "content_summary": "方法", "visual_prompt": "方法"},
                {"index": 4, "title": "价值", "content_summary": "价值", "visual_prompt": "价值"},
                {"index": 5, "title": "场景", "content_summary": "场景", "visual_prompt": "场景"},
                {"index": 6, "title": "风险", "content_summary": "风险", "visual_prompt": "风险"},
                {"index": 7, "title": "结论与建议", "content_summary": "建议", "visual_prompt": "建议"},
            ],
        },
    )
    monkeypatch.setattr(
        ppt_api.ppt_service.image_service,
        "generate_slide_image",
        lambda visual_prompt, reference_style_prompt=None: _build_test_png_bytes(),
    )

    response = client.post(
        "/ppt/generate",
        json={"document_ids": [document_id], "slide_count": 8},
    )
    assert response.status_code == 200
    outline = response.json()["data"]["outline"]
    titles = [slide["title"] for slide in outline]
    assert "背景" in titles
    assert "概念" in titles
    assert "方法" in titles
    assert "价值" in titles
    assert "风险" in titles
    assert "结论与行动建议" == titles[-1]


def test_generated_pptx_contains_single_full_slide_image_per_page(client, tmp_path, monkeypatch):
    configure_document_environment(tmp_path)

    monkeypatch.setattr(
        ppt_api.ppt_service.planner_service,
        "generate_outline",
        lambda topic, slide_count, context_text: {
            "global_style": "clean tech",
            "slides": [
                {
                    "index": index,
                    "title": f"Slide {index + 1}",
                    "content_summary": f"Summary {index + 1}",
                    "visual_prompt": f"Visual {index + 1}",
                }
                for index in range(slide_count)
            ],
        },
    )
    monkeypatch.setattr(
        ppt_api.ppt_service.image_service,
        "generate_slide_image",
        lambda visual_prompt, reference_style_prompt=None: _build_test_png_bytes(),
    )

    response = client.post("/ppt/generate", json={"topic": "AI Agent", "slide_count": 6})
    assert response.status_code == 200

    download_response = client.get(response.json()["data"]["download_url"].replace("/api/backend", ""))
    assert download_response.status_code == 200

    pptx_path = tmp_path / "deck.pptx"
    pptx_path.write_bytes(download_response.content)
    presentation = Presentation(pptx_path)
    assert len(presentation.slides) == 6
    assert presentation.slide_width == 12191999
    assert presentation.slide_height == 6858000
    for slide in presentation.slides:
        picture_shapes = [shape for shape in slide.shapes if shape.shape_type == 13]
        assert len(picture_shapes) == 1
        picture = picture_shapes[0]
        assert picture.left == 0
        assert picture.top == 0
        assert picture.width == presentation.slide_width
        assert picture.height == presentation.slide_height


def test_validate_slide_image_accepts_16_9_input(tmp_path):
    configure_document_environment(tmp_path)

    ppt_api.ppt_service._validate_slide_image(_build_test_png_bytes(size=(1600, 900), color=(180, 210, 240)))


def test_validate_slide_image_rejects_non_16_9_input(tmp_path):
    configure_document_environment(tmp_path)

    try:
        ppt_api.ppt_service._validate_slide_image(_build_test_png_bytes(size=(1024, 1024), color=(120, 160, 200)))
        raise AssertionError("Expected non-16:9 image to be rejected")
    except Exception as exc:
        assert "16:9" in str(exc)


def _build_test_png_bytes(
    size: tuple[int, int] = (1600, 900),
    color: tuple[int, int, int] = (248, 245, 239),
) -> bytes:
    from io import BytesIO
    from PIL import Image

    image = Image.new("RGB", size, color=color)
    buffer = BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()
